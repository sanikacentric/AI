import streamlit as st
import openai
import tempfile
import logging
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from PyPDF2 import PdfReader
import docx
import mammoth  # For .doc file support
from pptx import Presentation  # For .pptx support
import io

# Initialize the OpenAI API client
client = openai
logger = logging.getLogger(__name__)

# Extended Supported file types and their corresponding MIME types
SUPPORTED_FILE_TYPES = {
    'c': 'text/x-c',
    'cs': 'text/x-csharp',
    'cpp': 'text/x-c++',
    'doc': 'application/msword',
    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'html': 'text/html',
    'java': 'text/x-java',
    'json': 'application/json',
    'md': 'text/markdown',
    'pdf': 'application/pdf',
    'php': 'text/x-php',
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'py': 'text/x-python',
    'rb': 'text/x-ruby',
    'tex': 'text/x-tex',
    'txt': 'text/plain',
    'css': 'text/css',
    'js': 'text/javascript',
    'sh': 'application/x-sh',
    'ts': 'application/typescript',
    'csv': 'application/csv',
    'jpeg': 'image/jpeg',
    'jpg': 'image/jpeg',
    'gif': 'image/gif',
    'pkl': 'application/octet-stream',
    'png': 'image/png',
    'tar': 'application/x-tar',
    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'xml': 'application/xml',
    'zip': 'application/zip'
}

def read_file(uploaded_file, file_extension):
    """Read and process uploaded files based on their extension."""
    if file_extension == 'csv':
        return pd.read_csv(uploaded_file)
    elif file_extension == 'xlsx':
        return pd.read_excel(uploaded_file)
    elif file_extension == 'json':
        return pd.read_json(uploaded_file)
    elif file_extension == 'txt':
        return uploaded_file.read().decode('utf-8')
    elif file_extension == 'pdf':
        reader = PdfReader(uploaded_file)
        return "\n".join([page.extract_text() for page in reader.pages])
    elif file_extension == 'docx':
        doc = docx.Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs])
    elif file_extension == 'doc':
        with tempfile.NamedTemporaryFile(delete=True) as tmp:
            tmp.write(uploaded_file.read())
            tmp.flush()
            with open(tmp.name, "rb") as doc_file:
                result = mammoth.extract_raw_text(doc_file)
                return result.value
    elif file_extension == 'pptx':
        ppt = Presentation(uploaded_file)
        slides_text = []
        for slide in ppt.slides:
            slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            slides_text.append(slide_text)
        return "\n".join(slides_text)
    elif file_extension == 'xml':
        return uploaded_file.read().decode('utf-8')
    elif file_extension == 'html':
        return uploaded_file.read().decode('utf-8')
    elif file_extension in ['jpeg', 'jpg', 'png', 'gif']:
        st.image(uploaded_file)
        return None
    else:
        return None

def ask_openai_to_analyze(df, question):
    """Passes data to OpenAI for deeper insights and reasoning based on user query."""
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=".csv") as tmp:
        df.to_csv(tmp.name, index=False)
        tmp.seek(0)

        # Create an assistant with the uploaded file and user question
    try:
        assistant = client.beta.assistants.create(
            instructions=f"Analyze the dataset and answer: '{question}'",
            model="gpt-4",
            tools=[{"type": "code_interpreter"}],
            tool_resources={"code_interpreter": {"file_ids": [tmp.name]}}
        )
        st.success(f"Assistant created with ID: {assistant.id}")

        # Fetch the response and parse it
        response = assistant.response['choices'][0]['message']['content']
        st.write("Assistant Response:", response)

        # Display the response as a graph using Matplotlib
        if 'graph' in question.lower():
            plt.figure(figsize=(10, 6))
            sns.barplot(data=df)  # Example plotting using the dataframe
            plt.title("Assistant Suggested Graph")
            st.pyplot(plt)

    except Exception as e:
        st.error(f"Error during assistant response retrieval: {e}")
        logger.error(f"Error during assistant response retrieval: {e}")

def data_analysis_assistant():
    st.title("Data Analysis Assistant")

    # File upload widget, supporting various file formats
    uploaded_file = st.file_uploader("Upload your file", type=list(SUPPORTED_FILE_TYPES.keys()))

    if uploaded_file is not None:
        file_name = uploaded_file.name
        file_extension = file_name.split('.')[-1]

        file_content = read_file(uploaded_file, file_extension)
        
        # User input to ask questions about the data
        question = st.text_input("Ask a question about the data:")

        if st.button("Ask OpenAI"):
            ask_openai_to_analyze(file_content, question)

if __name__ == "__main__":
    data_analysis_assistant()
