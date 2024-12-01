import streamlit as st
import openai
import tempfile
import time
import logging
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from PyPDF2 import PdfReader
import docx
import mammoth  # For .doc file support
from pptx import Presentation  # For .pptx support
import io

# Initialize the OpenAI API client
client = openai
logger = logging.getLogger(__name__)

# Extended Supported file types and their corresponding MIME types
SUPPORTED_FILE_TYPES = {
    'c': 'text/x-c',
    'cs': 'text/x-csharp',
    'cpp': 'text/x-c++',
    'doc': 'application/msword',
    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'html': 'text/html',
    'java': 'text/x-java',
    'json': 'application/json',
    'md': 'text/markdown',
    'pdf': 'application/pdf',
    'php': 'text/x-php',
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'py': 'text/x-python',
    'rb': 'text/x-ruby',
    'tex': 'text/x-tex',
    'txt': 'text/plain',
    'css': 'text/css',
    'js': 'text/javascript',
    'sh': 'application/x-sh',
    'ts': 'application/typescript',
    'csv': 'application/csv',
    'jpeg': 'image/jpeg',
    'jpg': 'image/jpeg',
    'gif': 'image/gif',
    'pkl': 'application/octet-stream',
    'png': 'image/png',
    'tar': 'application/x-tar',
    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'xml': 'application/xml',
    'zip': 'application/zip'
}

def read_file(uploaded_file, file_extension):
    """Read and process uploaded files based on their extension."""
    if file_extension == 'csv':
        return pd.read_csv(uploaded_file)
    elif file_extension == 'xlsx':
        return pd.read_excel(uploaded_file)
    elif file_extension == 'json':
        return pd.read_json(uploaded_file)
    elif file_extension == 'txt':
        return uploaded_file.read().decode('utf-8')
    elif file_extension == 'pdf':
        reader = PdfReader(uploaded_file)
        return "\n".join([page.extract_text() for page in reader.pages])
    elif file_extension == 'docx':
        doc = docx.Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs])
    elif file_extension == 'doc':
        # Using mammoth to extract text from .doc files
        with tempfile.NamedTemporaryFile(delete=True) as tmp:
            tmp.write(uploaded_file.read())
            tmp.flush()
            with open(tmp.name, "rb") as doc_file:
                result = mammoth.extract_raw_text(doc_file)
                return result.value
    elif file_extension == 'pptx':
        # Using python-pptx to extract text from presentation slides
        ppt = Presentation(uploaded_file)
        slides_text = []
        for slide in ppt.slides:
            slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            slides_text.append(slide_text)
        return "\n".join(slides_text)
    elif file_extension == 'xml':
        return uploaded_file.read().decode('utf-8')
    elif file_extension == 'html':
        return uploaded_file.read().decode('utf-8')
    elif file_extension in ['jpeg', 'jpg', 'png', 'gif']:
        st.image(uploaded_file)
        return None
    else:
        return None

def analyze_data(df):
    """Performs in-depth data analysis on the DataFrame and returns insights."""
    st.write("### Data Summary")
    st.write(df.describe())

    # Identifying correlations between numeric columns
    if len(df.select_dtypes(include=['number']).columns) > 1:
        st.write("### Correlation Heatmap")
        corr = df.corr()
        sns.heatmap(corr, annot=True, cmap='coolwarm', linewidths=0.5)
        st.pyplot(plt)
    
    # Detecting missing values
    st.write("### Missing Data")
    missing_values = df.isnull().sum()
    if missing_values.sum() > 0:
        st.write(missing_values[missing_values > 0])
    else:
        st.write("No missing values detected.")

def generate_dynamic_charts(df):
    """Generates charts dynamically by detecting appropriate columns for x and y axes."""
    st.write("### Dynamic Bar Chart based on Data")
    
    # Identify potential columns for x-axis (categorical) and y-axis (numeric)
    potential_x_cols = df.select_dtypes(include=['object', 'category']).columns
    potential_y_cols = df.select_dtypes(include=['number']).columns

    if not potential_x_cols.empty and not potential_y_cols.empty:
        # Use the first detected categorical column for x-axis and numeric for y-axis
        x_col = potential_x_cols[0]
        y_col = potential_y_cols[0]
        
        st.write(f"Automatically selected `{x_col}` for the x-axis and `{y_col}` for the y-axis.")

        plt.figure(figsize=(12, 6))
        sns.barplot(x=df[x_col], y=df[y_col])
        plt.xticks(rotation=90)
        plt.title(f"Bar Chart of {x_col} vs {y_col}")
        plt.xlabel(x_col)
        plt.ylabel(y_col)
        st.pyplot(plt)
    else:
        st.write("Error: Unable to find appropriate columns for generating a chart.")


def ask_openai_to_analyze(df, question):
    """Passes data to OpenAI for deeper insights and reasoning based on user query."""
    # Upload the CSV content to OpenAI
    csv_buffer = io.StringIO()
    df.to_csv(csv_buffer)
    csv_content = csv_buffer.getvalue()

    with tempfile.NamedTemporaryFile(delete=False, suffix=".csv") as tmp:
        tmp.write(csv_content.encode('utf-8'))
        tmp.seek(0)

        # Upload the file to OpenAI with the purpose 'assistants'
        file = client.files.create(file=open(tmp.name, "rb"), purpose='assistants')
        st.success(f"File uploaded to OpenAI with ID: {file.id}")

        # Create an assistant with the uploaded file and user question
    try:
        assistant = client.beta.assistants.create(
            instructions=f"Analyze the dataset and answer: '{question}'",
            model="gpt-4",
            tools=[{"type": "code_interpreter"}],
            tool_resources={"code_interpreter": {"file_ids": [file.id]}}
        )
        st.success(f"Assistant created with ID: {assistant.id}")

        # Polling for the assistant's response
        response_text = "Generate a bar chart of table volumes"  # For simplicity, using a predefined response
        generate_dynamic_charts(df)  # Updated to use the dynamic function

    except Exception as e:
        st.error(f"Error during assistant response retrieval: {e}")
        logger.error(f"Error during assistant response retrieval: {e}")

def data_analysis_assistant():
    st.title("Data Analysis Assistant")

    # File upload widget, supporting various file formats
    uploaded_file = st.file_uploader("Upload your file", type=list(SUPPORTED_FILE_TYPES.keys()))

    if uploaded_file is not None:
        try:
            file_name = uploaded_file.name
            file_extension = file_name.split('.')[-1]
            st.write(f"Uploaded file: {file_name}")

            # Process the uploaded file based on its extension
            file_content = read_file(uploaded_file, file_extension)
            if isinstance(file_content, pd.DataFrame):
                st.dataframe(file_content)

                # Perform in-depth data analysis
                analyze_data(file_content)

                # User input to ask questions about the data
                question = st.text_input("Ask a question about the data (e.g., 'Generate a bar chart of table volumes'):")

                if st.button("Ask OpenAI"):
                    if question and isinstance(file_content, pd.DataFrame):
                        ask_openai_to_analyze(file_content, question)

        except Exception as e:
            st.error(f"Error during file upload or analysis: {e}")
            logger.error(f"Error during file upload or analysis: {e}")

if __name__ == "__main__":
    data_analysis_assistant()