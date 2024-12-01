import streamlit as st
import openai
import tempfile
import time
import logging
import pandas as pd
import matplotlib.pyplot as plt
from PyPDF2 import PdfReader
import docx
import mammoth  # For .doc file support
from pptx import Presentation  # For .pptx support

# Initialize the OpenAI API client
client = openai
logger = logging.getLogger(__name__)

# Extended Supported file types and their corresponding MIME types
SUPPORTED_FILE_TYPES = {
    'c': 'text/x-c',
    'cs': 'text/x-csharp',
    'cpp': 'text/x-c++',
    'doc': 'application/msword',  # Added .doc support
    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'html': 'text/html',
    'java': 'text/x-java',
    'json': 'application/json',
    'md': 'text/markdown',
    'pdf': 'application/pdf',
    'php': 'text/x-php',
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # Added .pptx support
    'py': 'text/x-python',
    'rb': 'text/x-ruby',
    'tex': 'text/x-tex',
    'txt': 'text/plain',
    'css': 'text/css',
    'js': 'text/javascript',
    'sh': 'application/x-sh',
    'ts': 'application/typescript',
    'csv': 'application/csv',
    'jpeg': 'image/jpeg',
    'jpg': 'image/jpeg',
    'gif': 'image/gif',
    'pkl': 'application/octet-stream',
    'png': 'image/png',
    'tar': 'application/x-tar',
    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'xml': 'application/xml',
    'zip': 'application/zip'
}

def read_file(uploaded_file, file_extension):
    """Read and process uploaded files based on their extension."""
    if file_extension == 'csv':
        return pd.read_csv(uploaded_file)
    elif file_extension == 'xlsx':
        return pd.read_excel(uploaded_file)
    elif file_extension == 'json':
        return pd.read_json(uploaded_file)
    elif file_extension == 'txt':
        return uploaded_file.read().decode('utf-8')
    elif file_extension == 'pdf':
        reader = PdfReader(uploaded_file)
        return "\n".join([page.extract_text() for page in reader.pages])
    elif file_extension == 'docx':
        doc = docx.Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs])
    elif file_extension == 'doc':
        # Using mammoth to extract text from .doc files
        with tempfile.NamedTemporaryFile(delete=True) as tmp:
            tmp.write(uploaded_file.read())
            tmp.flush()
            with open(tmp.name, "rb") as doc_file:
                result = mammoth.extract_raw_text(doc_file)
                return result.value
    elif file_extension == 'pptx':
        # Using python-pptx to extract text from presentation slides
        ppt = Presentation(uploaded_file)
        slides_text = []
        for slide in ppt.slides:
            slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            slides_text.append(slide_text)
        return "\n".join(slides_text)
    elif file_extension == 'xml':
        return uploaded_file.read().decode('utf-8')
    elif file_extension == 'html':
        return uploaded_file.read().decode('utf-8')
    elif file_extension in ['jpeg', 'jpg', 'png', 'gif']:
        st.image(uploaded_file)
        return None
    else:
        return None

def data_analysis_assistant():
    st.title("Data Analysis Assistant")

    # File upload widget, supporting various file formats
    uploaded_file = st.file_uploader(
        "Upload your file", 
        type=list(SUPPORTED_FILE_TYPES.keys())
    )

    if uploaded_file is not None:
        try:
            file_name = uploaded_file.name
            file_extension = file_name.split('.')[-1]

            st.write(f"Uploaded file: {file_name}")

            # Process the uploaded file based on its extension
            file_content = read_file(uploaded_file, file_extension)
            if file_content is not None:
                if isinstance(file_content, pd.DataFrame):
                    st.dataframe(file_content)
                else:
                    st.text_area("File Content", file_content)

            # User input to ask questions about the data
            question = st.text_input("Ask a question about the data (e.g., 'Show summary' or 'Generate a bar chart of column X'):")

            if question:
                # Example responses based on the input question
                if isinstance(file_content, pd.DataFrame) and 'summary' in question.lower():
                    st.write("Summary of the dataset:")
                    st.write(file_content.describe())
                elif isinstance(file_content, pd.DataFrame) and 'bar chart' in question.lower():
                    # Extract the column name from the question
                    column_name = question.split("of")[-1].strip()
                    if column_name in file_content.columns:
                        st.write(f"Generating bar chart for '{column_name}'")
                        plt.figure(figsize=(10, 6))
                        file_content[column_name].value_counts().plot(kind='bar')
                        st.pyplot(plt)
                    else:
                        st.error(f"Column '{column_name}' not found in the dataset.")
                else:
                    # Pass the question to the OpenAI assistant for complex analysis
                    file_data = uploaded_file.read()

                    # Handle file upload securely with tempfile
                    with tempfile.NamedTemporaryFile(delete=True, suffix=f".{file_extension}") as tmp:
                        tmp.write(file_data)
                        tmp.seek(0)

                        # Upload the file to OpenAI with the purpose 'assistants'
                        file = client.files.create(
                            file=open(tmp.name, "rb"),
                            purpose='assistants'
                        )
                    
                    st.success(f"File '{file_name}' uploaded to assistant with ID: {file.id}")

                    # Create an assistant with the uploaded file and enable Code Interpreter
                    assistant = client.beta.assistants.create(
                        instructions=f"You are a data analyst. Please answer the following question based on the uploaded file: '{question}'",
                        model="gpt-4",  # Replace with your required model
                        tools=[{"type": "code_interpreter"}],
                        tool_resources={
                            "code_interpreter": {
                                "file_ids": [file.id]
                            }
                        }
                    )

                    st.success(f"Assistant created successfully with ID: {assistant.id}")

                    # Create a thread and pass the file and question as part of the message creation request
                    thread = client.beta.threads.create(
                        messages=[{
                            "role": "user",
                            "content": question,
                            "attachments": [{
                                "file_id": file.id,
                                "tools": [{"type": "code_interpreter"}]
                            }]
                        }]
                    )

                    st.success(f"Thread created successfully with ID: {thread.id}")

                    # Polling or waiting for the assistant's message response
                    if st.button("Analyze Data"):
                        max_retries = 20  # Increased from 10 to 20
                        retry_count = 0
                        message_found = False

                        while retry_count < max_retries:
                            retry_count += 1
                            st.write(f"Checking for message response... Attempt {retry_count}/{max_retries}")

                            # Retrieve messages from the thread to check if the assistant responded
                            thread_messages = client.beta.threads.messages.list(thread_id=thread.id)
                            st.write(f"Retrieved Messages (Attempt {retry_count}): {thread_messages.data}")

                            for message in thread_messages.data:
                                if message.role == 'assistant':
                                    st.write(f"Message from Assistant: {message.content}")
                                    message_found = True
                                    break

                            if message_found:
                                break
                            time.sleep(10)  # Increased sleep duration to 10 seconds between attempts

                        if not message_found:
                            st.write("No response from the assistant yet. Please try again later.")

        except Exception as e:
            st.error(f"Error during file upload or analysis: {e}")
            logger.error(f"Error during file upload or analysis: {e}")
